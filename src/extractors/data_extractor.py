import os
import csv
import io
import json

import fitz        # For PDF
import pdfplumber  # Also for PDF table extraction
from docx import Document
from pptx import Presentation
from PIL import Image

class PDFDataExtractor:
    """
    Extracts text, links, images, and tables from a PyMuPDF Document object.
    """

    def __init__(self, pdf_doc, file_path: str):
        self.pdf_doc = pdf_doc     # a fitz Document
        self.file_path = file_path # so we can name output folders, etc.

    def extract_text(self):
        text_content = {}
        font_styles = []
        headings = {}

        for page_index, page in enumerate(self.pdf_doc):
            page_num = page_index + 1
            blocks = page.get_text("dict")["blocks"]
            page_text = []

            for block in blocks:
                for line in block.get("lines", []):
                    for span in line.get("spans", []):
                        font_size = span["size"]
                        font_name = span["font"]
                        extracted_text = span["text"].strip()

                        if extracted_text:
                            page_text.append(extracted_text)
                            font_styles.append({
                                "page_number": page_num,
                                "text": extracted_text,
                                "font": font_name,
                                "size": font_size
                            })
                            # Simple heading heuristic
                            if font_size >= 14:
                                if page_num not in headings:
                                    headings[page_num] = []
                                headings[page_num].append(extracted_text)

            text_content[page_num] = page_text

        return {
            "text": text_content,
            "metadata": {
                "font_styles": font_styles,
                "headings": headings
            }
        }

    def extract_links(self):
        links_info = []
        for page_index, page in enumerate(self.pdf_doc):
            page_num = page_index + 1
            links = page.get_links()
            for link in links:
                if 'uri' in link:
                    links_info.append({
                        "url": link['uri'],
                        "page_number": page_num
                    })
        return links_info

    def extract_images(self):
        images_info = []
        base_name = os.path.splitext(os.path.basename(self.file_path))[0]
        output_dir = os.path.join("output", base_name, "images")
        os.makedirs(output_dir, exist_ok=True)

        for page_index, page in enumerate(self.pdf_doc):
            page_num = page_index + 1
            image_list = page.get_images(full=True)
            for img_index, img in enumerate(image_list):
                xref = img[0]
                base_image = self.pdf_doc.extract_image(xref)
                image_data = base_image["image"]
                ext = base_image["ext"]

                img_filename = f"page_{page_num}_img_{img_index}.{ext}"
                img_path = os.path.join(output_dir, img_filename)

                with open(img_path, "wb") as f:
                    f.write(image_data)

                images_info.append({
                    "page_number": page_num,
                    "image_path": img_path
                })
        return images_info

    def extract_tables(self):
        tables_info = []
        base_name = os.path.splitext(os.path.basename(self.file_path))[0]
        output_dir = os.path.join("output/tables", base_name)
        os.makedirs(output_dir, exist_ok=True)

        # Use pdfplumber for table extraction
        with pdfplumber.open(self.file_path) as pdf_file:
            for page_index, page in enumerate(pdf_file.pages):
                page_num = page_index + 1
                extracted_table = page.extract_table()
                if extracted_table:
                    table_filename = f"page_{page_num}_table.csv"
                    table_path = os.path.join(output_dir, table_filename)
                    with open(table_path, "w", newline="", encoding="utf-8") as f:
                        writer = csv.writer(f)
                        for row in extracted_table:
                            writer.writerow(row)

                    tables_info.append({
                        "page_number": page_num,
                        "table_data": extracted_table,
                        "table_path": table_path
                    })

        return tables_info


class DOCXDataExtractor:
    """
    Extracts text, links, images, and tables from a python-docx Document object.
    """

    def __init__(self, docx_doc, file_path: str):
        self.docx_doc = docx_doc   # a Document
        self.file_path = file_path

    def extract_text(self):
        text_content = {}
        headings = {}
        font_styles = []
        page_num = 1  # For DOCX we typically treat everything as page 1

        all_paragraphs = []
        for para in self.docx_doc.paragraphs:
            txt = para.text.strip()
            if txt:
                all_paragraphs.append(txt)
                style_name = para.style.name if para.style else "Normal"
                font_styles.append({
                    "page_number": page_num,
                    "text": txt,
                    "font": style_name,
                    "size": 12  # rough default
                })
                if style_name.startswith("Heading"):
                    headings.setdefault(page_num, []).append(txt)

        text_content[page_num] = all_paragraphs
        return {
            "text": text_content,
            "metadata": {
                "headings": headings,
                "font_styles": font_styles
            }
        }

    def extract_links(self):
        links_info = []
        page_num = 1

        for para in self.docx_doc.paragraphs:
            for run in para.runs:
                if hasattr(run._element, 'xpath'):
                    hyperlink_refs = run._element.xpath('.//w:hyperlink')
                    for link_ref in hyperlink_refs:
                        rel_id = link_ref.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
                        if rel_id and rel_id in self.docx_doc.part.rels:
                            target = self.docx_doc.part.rels[rel_id].target_ref
                            link_text = run.text.strip()
                            links_info.append({
                                "url": target,
                                "page_number": page_num,
                                "text": link_text
                            })

        return links_info

    def extract_images(self):
        images_info = []
        base_name = os.path.splitext(os.path.basename(self.file_path))[0]
        output_dir = os.path.join("output", base_name, "images")
        os.makedirs(output_dir, exist_ok=True)

        image_index = 0
        for rel in self.docx_doc.part.rels.values():
            if "image" in rel.reltype:
                img_ext = os.path.splitext(rel.target_ref)[1]
                img_path = os.path.join(output_dir, f"image_{image_index}{img_ext}")
                with open(img_path, "wb") as f:
                    f.write(rel.target_part.blob)
                images_info.append({
                    "page_number": 1,
                    "image_path": img_path
                })
                image_index += 1
        return images_info

    def extract_tables(self):
        tables_info = []
        base_name = os.path.splitext(os.path.basename(self.file_path))[0]
        output_dir = os.path.join("output/tables", base_name)
        os.makedirs(output_dir, exist_ok=True)

        for tbl_index, table in enumerate(self.docx_doc.tables):
            extracted_data = []
            for row in table.rows:
                row_data = [cell.text.strip() for cell in row.cells]
                extracted_data.append(row_data)

            table_filename = f"table_{tbl_index}.csv"
            table_path = os.path.join(output_dir, table_filename)
            with open(table_path, "w", newline="", encoding="utf-8") as f:
                writer = csv.writer(f)
                writer.writerows(extracted_data)

            tables_info.append({
                "page_number": 1,
                "table_index": tbl_index,
                "table_data": extracted_data,
                "table_path": table_path
            })

        return tables_info

class PPTDataExtractor:
    """
    Extracts text, links, images, and tables from a python-pptx Presentation object.
    """

    def __init__(self, ppt_doc, file_path: str):
        self.ppt_doc = ppt_doc
        self.file_path = file_path

    def extract_text(self):
        text_content = {}
        headings = {}
        font_styles = []

        for i, slide in enumerate(self.ppt_doc.slides):
            slide_num = i + 1
            slide_text_list = []

            # Title placeholder
            if slide.shapes.title and hasattr(slide.shapes.title, 'text'):
                title_text = slide.shapes.title.text.strip()
                if title_text:
                    slide_text_list.append(title_text)
                    headings.setdefault(slide_num, []).append(title_text)

            # All shapes
            for shape in slide.shapes:
                if hasattr(shape, "text_frame") and shape.text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        paragraph_text = paragraph.text.strip()
                        if paragraph_text:
                            slide_text_list.append(paragraph_text)
                            # run-level font info
                            for run in paragraph.runs:
                                run_txt = run.text.strip()
                                if run_txt:
                                    size_pts = run.font.size.pt if run.font.size else 12
                                    font_styles.append({
                                        "page_number": slide_num,
                                        "text": run_txt,
                                        "font": run.font.name if run.font.name else "Default",
                                        "size": size_pts
                                    })

            text_content[slide_num] = slide_text_list

        return {
            "text": text_content,
            "metadata": {
                "headings": headings,
                "font_styles": font_styles
            }
        }

    def extract_links(self):
        links_info = []
        for i, slide in enumerate(self.ppt_doc.slides):
            slide_num = i + 1
            for shape in slide.shapes:
                # 1) Shape-level hyperlink (click_action)
                if hasattr(shape, "click_action") and shape.click_action:
                    if (hasattr(shape.click_action, "hyperlink") 
                            and shape.click_action.hyperlink 
                            and hasattr(shape.click_action.hyperlink, "address") 
                            and shape.click_action.hyperlink.address):
                        # Grab shape text from text_frame, if present
                        shape_text = ""
                        if hasattr(shape, "text_frame") and shape.text_frame:
                            shape_text = shape.text_frame.text.strip()
                        links_info.append({
                            "page_number": slide_num,
                            "text": shape_text,
                            "url": shape.click_action.hyperlink.address,
                            "shape_name": shape.name if hasattr(shape, "name") else ""
                        })

                # 2) Run-level hyperlinks (within text_frame paragraphs)
                if hasattr(shape, "text_frame") and shape.text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if run.hyperlink and run.hyperlink.address:
                                links_info.append({
                                    "page_number": slide_num,
                                    "text": run.text.strip(),
                                    "url": run.hyperlink.address,
                                    "shape_name": shape.name if hasattr(shape, "name") else ""
                                })
        return links_info

    def extract_images(self):
        images_info = []
        base_name = os.path.splitext(os.path.basename(self.file_path))[0]
        output_dir = os.path.join("output", base_name, "images")
        os.makedirs(output_dir, exist_ok=True)

        image_index = 0
        for slide_index, slide in enumerate(self.ppt_doc.slides):
            slide_num = slide_index + 1
            for shape in slide.shapes:
                # shape_type == 13 => PICTURE
                if shape.shape_type == 13 and hasattr(shape, "image"):
                    try:
                        img_stream = io.BytesIO(shape.image.blob)
                        pil_img = Image.open(img_stream)
                        img_filename = f"slide_{slide_num}_img_{image_index}.png"
                        out_path = os.path.join(output_dir, img_filename)
                        pil_img.save(out_path)

                        images_info.append({
                            "page_number": slide_num,
                            "image_path": out_path,
                            # Optional: capture alt_text if present
                            "alt_text": shape.alt_text if hasattr(shape, "alt_text") else ""
                        })
                        image_index += 1
                    except Exception as e:
                        print(f"Error extracting image on slide {slide_num}: {e}")
        return images_info

    def extract_tables(self):
        tables_info = []
        base_name = os.path.splitext(os.path.basename(self.file_path))[0]
        output_dir = os.path.join("output/tables", base_name)
        os.makedirs(output_dir, exist_ok=True)

        table_count = 0
        for slide_index, slide in enumerate(self.ppt_doc.slides):
            slide_num = slide_index + 1
            for shape in slide.shapes:
                if hasattr(shape, "has_table") and shape.has_table:
                    table_data = []
                    for row in shape.table.rows:
                        row_data = [cell.text.strip() for cell in row.cells]
                        table_data.append(row_data)

                    csv_filename = f"slide_{slide_num}_table_{table_count}.csv"
                    csv_path = os.path.join(output_dir, csv_filename)
                    with open(csv_path, "w", newline="", encoding="utf-8") as f:
                        writer = csv.writer(f)
                        writer.writerows(table_data)

                    tables_info.append({
                        "page_number": slide_num,
                        "table_data": table_data,
                        "table_path": csv_path
                    })
                    table_count += 1

        return tables_info
